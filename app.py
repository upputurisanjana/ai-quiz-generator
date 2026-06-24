import streamlit as st
import streamlit.components.v1 as components
from pptx import Presentation
import os, json, time, math
from dotenv import load_dotenv

load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), ".env"))
st.set_page_config(page_title="Scholarium", page_icon="📜", layout="wide")

_DEFAULTS = {
    "screen": "upload", "ppt_data": None,
    "quiz_config": {}, "questions": [], "current_question": 0,
    "user_answers": {}, "confirmed_answers": {}, "timer_start": None,
}
for k, v in _DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v

def inject_css():
    st.markdown("""
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,400;0,500;0,600;1,400;1,500&family=Crimson+Pro:ital,wght@0,400;0,600;1,400&family=Cinzel:wght@400;500;600&display=swap" rel="stylesheet">
<style>
:root {
  --bg:     #1C1714;
  --bg-alt: #251E19;
  --fg:     #E8DFD4;
  --muted:  #3D332B;
  --mfg:    #9C8B7A;
  --bd:     #4A3F35;
  --ac:     #C9A962;
  --ac2:    #8B2635;
  --ac-fg:  #1C1714;
  --brass-grad: linear-gradient(180deg,#D4B872 0%,#C9A962 50%,#B8953F 100%);
  --engraved: 1px 1px 1px rgba(0,0,0,.4),-1px -1px 1px rgba(255,255,255,.1);
  --fh: 'Cormorant Garamond',serif;
  --fb: 'Crimson Pro',serif;
  --fd: 'Cinzel',serif;
}
html,body,.stApp { background:var(--bg)!important; color:var(--fg); font-family:var(--fb); }
.stApp > header, .stDeployButton, footer,
.st-emotion-cache-18ni7ap, .st-emotion-cache-6qob1r { display:none!important; }
section[data-testid="stSidebar"] { display:none!important; }
.main .block-container { padding:52px 60px 100px!important; max-width:100%!important; min-height:100vh; }
.ac-texture {
  position:fixed; inset:0; pointer-events:none; z-index:0;
  background-image: url("data:image/svg+xml,%3Csvg xmlns=%22http://www.w3.org/2000/svg%22 width=%22300%22 height=%22300%22%3E%3Cfilter id=%22n%22%3E%3CfeTurbulence type=%22fractalNoise%22 baseFrequency=%220.75%22 numOctaves=%224%22 stitchTiles=%22stitch%22/%3E%3C/filter%3E%3Crect width=%22300%22 height=%22300%22 filter=%22url(%23n)%22 opacity=%220.03%22/%3E%3C/svg%3E");
  mix-blend-mode:overlay;
}
.ac-vignette {
  position:fixed; inset:0; pointer-events:none; z-index:0;
  background: radial-gradient(ellipse at center, transparent 0%, transparent 50%, rgba(28,23,20,.5) 100%);
}
.ac-overline { font-family:var(--fd); font-size:10px; font-weight:500; text-transform:uppercase; letter-spacing:.25em; color:var(--ac); }
.ac-h1 { font-family:var(--fh); font-size:52px; font-weight:400; line-height:1.05; letter-spacing:-.01em; color:var(--fg); margin:8px 0 14px; }
.ac-sub { font-family:var(--fb); font-size:17px; color:var(--mfg); line-height:1.6; margin:0 0 36px; }
.ac-div { position:relative; height:1px; margin:24px 0; background:linear-gradient(90deg,transparent 0%,var(--bd) 20%,var(--ac) 50%,var(--bd) 80%,transparent 100%); }
.ac-div::before { content:"✶"; position:absolute; left:50%; top:50%; transform:translate(-50%,-50%); color:var(--ac); font-size:11px; background:var(--bg-alt); padding:0 10px; }
.ac-card { background:var(--bg-alt); border:1px solid var(--bd); border-radius:4px; padding:28px 32px; position:relative; overflow:hidden; transition:border-color .3s ease, box-shadow .3s ease; }
.ac-card:hover { border-color:rgba(201,169,98,.45); box-shadow:0 8px 24px rgba(0,0,0,.3); }
.ac-card::before, .ac-card::after { content:""; position:absolute; width:24px; height:24px; border:1.5px solid var(--ac); opacity:.5; transition:opacity .3s; }
.ac-card::before { top:0; left:0; border-right:none; border-bottom:none; }
.ac-card::after  { bottom:0; right:0; border-left:none; border-top:none; }
.ac-card:hover::before, .ac-card:hover::after { opacity:1; }
.ac-inset { background:var(--bg); border:1px solid var(--bd); border-radius:4px; padding:22px 24px; }
div.stButton > button { background:var(--brass-grad)!important; color:var(--ac-fg)!important; border:none!important; border-radius:4px!important; font-family:var(--fd)!important; font-size:11px!important; font-weight:500!important; text-transform:uppercase!important; letter-spacing:.15em!important; padding:14px 28px!important; text-shadow:var(--engraved)!important; box-shadow:inset 0 1px 0 rgba(255,255,255,.2),inset 0 -1px 0 rgba(0,0,0,.2),0 2px 8px rgba(0,0,0,.3)!important; transition:filter .15s ease,box-shadow .15s ease!important; }
div.stButton > button:hover { filter:brightness(1.1)!important; box-shadow:inset 0 1px 0 rgba(255,255,255,.2),inset 0 -1px 0 rgba(0,0,0,.2),0 4px 12px rgba(201,169,98,.3)!important; }
div.stButton > button:active { box-shadow:inset 0 2px 4px rgba(0,0,0,.35)!important; }
div.stButton > button:focus-visible { outline:none!important; box-shadow:0 0 0 2px var(--bg),0 0 0 4px var(--ac)!important; }
div[data-testid="stSlider"] > div { padding:0!important; }
div[data-testid="stSlider"] label { font-family:var(--fd)!important; font-size:10px!important; font-weight:500!important; text-transform:uppercase!important; letter-spacing:.2em!important; color:var(--mfg)!important; }
div[data-testid="stSlider"] [data-baseweb="slider"] [role="slider"] { background:var(--ac)!important; border-color:var(--ac)!important; }
div[data-testid="stSlider"] [data-baseweb="slider"] [data-testid="stThumbValue"] { font-family:var(--fd)!important; color:var(--ac)!important; font-size:11px!important; }
div[data-testid="stExpander"] { background:var(--bg-alt)!important; border:1px solid var(--bd)!important; border-radius:4px!important; margin-bottom:6px!important; }
div[data-testid="stExpander"] summary { font-family:var(--fh)!important; font-size:16px!important; color:var(--fg)!important; padding:14px 18px!important; }
.ac-prog { height:2px; background:var(--muted); border-radius:99px; overflow:hidden; margin-bottom:28px; }
.ac-progf { height:100%; background:var(--brass-grad); border-radius:99px; transition:width .5s ease-out; }
.ac-opt { width:100%; background:var(--bg-alt); border:1px solid var(--bd); border-radius:4px; padding:14px 18px; text-align:left; cursor:pointer; font-family:var(--fb); font-size:16px; color:var(--fg); display:flex; align-items:center; gap:14px; margin-bottom:8px; transition:border-color .15s, background .15s; }
.ac-opt:hover { border-color:rgba(201,169,98,.5); }
.ac-opt.sel { border-color:var(--ac); background:rgba(201,169,98,.08); }
.ac-opt.ok  { border-color:#4A7C5F; background:rgba(74,124,95,.1); }
.ac-opt.err { border-color:#8B3030; background:rgba(139,48,48,.1); opacity:.85; }
.ac-key { font-family:var(--fd); font-size:10px; font-weight:500; text-transform:uppercase; letter-spacing:.1em; background:var(--muted); color:var(--mfg); border-radius:2px; padding:2px 7px; flex-shrink:0; transition:background .15s,color .15s; }
.ac-opt.sel .ac-key { background:var(--ac); color:var(--ac-fg); }
.ac-opt.ok  .ac-key { background:#4A7C5F; color:#E8DFD4; }
.ac-opt.err .ac-key { background:#8B3030; color:#E8DFD4; }
.ac-xrow { padding:14px 18px; font-family:var(--fb); font-size:15px; line-height:1.6; border-left:2px solid; margin-bottom:5px; border-radius:0 4px 4px 0; }
.xr-ok  { border-color:#4A7C5F; background:rgba(74,124,95,.08); }
.xr-err { border-color:#8B3030; background:rgba(139,48,48,.08); }
.ac-xlbl { font-family:var(--fd); font-size:9px; font-weight:500; text-transform:uppercase; letter-spacing:.2em; margin-bottom:6px; }
.xr-ok  .ac-xlbl { color:#7AAF95; }
.xr-err .ac-xlbl { color:#B87070; }
.ac-tw { position:relative; width:72px; height:72px; margin:0 auto 8px; }
.ac-tn { position:absolute; inset:0; display:flex; align-items:center; justify-content:center; font-family:var(--fh); font-size:20px; font-weight:500; color:var(--fg); transition:color .3s; }
.ac-tn.urg { color:#B87070; animation:pulse 1s ease infinite; }
.ac-score { font-family:var(--fh); font-size:72px; font-weight:400; letter-spacing:-.02em; line-height:1; color:var(--ac); }
.ac-score-denom { font-family:var(--fh); font-size:36px; color:var(--mfg); }
.ac-seal { display:inline-flex; align-items:center; justify-content:center; width:36px; height:36px; border-radius:50%; background:radial-gradient(circle at 40% 35%,#A83248 0%,#8B2635 55%,#6B1A28 100%); box-shadow:inset 0 2px 4px rgba(255,255,255,.2),inset 0 -2px 4px rgba(0,0,0,.3),0 3px 7px rgba(0,0,0,.4); color:#E8DFD4; font-size:14px; flex-shrink:0; }
.ac-drop { border:1px solid var(--bd); border-radius:4px; padding:56px 36px; text-align:center; position:relative; overflow:hidden; background:var(--bg-alt); transition:border-color .3s; }
.ac-drop:hover { border-color:rgba(201,169,98,.45); }
.ac-arch { width:80px; height:80px; margin:0 auto 20px; border:1px solid var(--bd); border-radius:40% 40% 0 0 / 20% 20% 0 0; display:flex; align-items:center; justify-content:center; background:var(--bg); color:var(--ac); font-size:28px; transition:border-color .3s; }
.ac-drop:hover .ac-arch { border-color:rgba(201,169,98,.6); }
.ac-lbar { width:100%; height:2px; background:var(--muted); border-radius:99px; overflow:hidden; }
.ac-lbarf { height:100%; background:var(--brass-grad); border-radius:99px; animation:lsc 1.6s ease-in-out infinite; }
div[data-testid^="qngrid_"] button { background:var(--bg-alt)!important; color:var(--mfg)!important; border:1px solid var(--bd)!important; border-radius:4px!important; padding:0!important; min-height:34px!important; font-family:var(--fd)!important; font-size:10px!important; font-weight:500!important; }
div[data-testid^="qngrid_"] button:hover { border-color:var(--ac)!important; color:var(--ac)!important; }
div[data-testid^="opt_btn_"] { position:relative!important; margin-bottom:8px!important; }
div[data-testid^="opt_btn_"] div[data-testid="stButton"] { position:absolute!important; inset:0!important; margin:0!important; }
div[data-testid^="opt_btn_"] button { position:absolute!important; inset:0!important; width:100%!important; height:100%!important; opacity:0!important; cursor:pointer!important; border:none!important; background:none!important; padding:0!important; box-shadow:none!important; }
.ac-chip { display:inline-flex; align-items:center; gap:8px; background:rgba(201,169,98,.1); border:1px solid rgba(201,169,98,.3); border-radius:2px; padding:4px 14px; font-family:var(--fd); font-size:10px; font-weight:500; text-transform:uppercase; letter-spacing:.15em; color:var(--ac); margin-bottom:24px; }
@keyframes lsc { 0%{width:0;margin-left:0;} 50%{width:55%;margin-left:22%;} 100%{width:0;margin-left:100%;} }
@keyframes pulse { 0%,100%{opacity:1;} 50%{opacity:.45;} }
@keyframes fup { from{opacity:0;transform:translateY(8px);} to{opacity:1;transform:translateY(0);} }
.ac-anim { animation:fup .3s ease-out; }
@media(prefers-reduced-motion:reduce) { *,*::before,*::after { animation-duration:.01ms!important; transition-duration:.01ms!important; } }
</style>
<div class="ac-texture" aria-hidden="true"></div>
<div class="ac-vignette" aria-hidden="true"></div>
""", unsafe_allow_html=True)

_TIME_PER_Q = {"Simple": 35, "Medium": 67, "Complex": 105}

def parse_pptx(f):
    prs = Presentation(f)
    slides = [{"n": i+1, "t": " ".join(sh.text for sh in s.shapes if hasattr(sh,"text"))} for i,s in enumerate(prs.slides)]
    return {"filename": f.name, "slides": len(slides), "words": sum(len(s["t"].split()) for s in slides), "raw": slides}

def generate_questions(ppt, n, diff):
    from openai import OpenAI
    key = os.environ.get("OPENROUTER_API_KEY")
    if not key:
        st.error("OPENROUTER_API_KEY missing"); st.stop()
    client = OpenAI(api_key=key, base_url="https://openrouter.ai/api/v1")
    txt = "\n\n".join(f"Slide {s['n']}: {s['t']}" for s in ppt["raw"])
    prompt = f"""Generate {n} MCQs at {diff} difficulty.
Return ONLY a JSON array:
[{{"id":"q1","question_text":"...","options":{{"A":"...","B":"...","C":"...","D":"..."}},"correct_answer":"A","source_slide":1,"explanation_correct":"...","explanation_wrong":{{"B":"...","C":"...","D":"..."}}}}]
Content:\n{txt}"""
    r = client.chat.completions.create(model="deepseek/deepseek-r1", max_tokens=4096,
        messages=[{"role":"user","content":prompt}])
    c = r.choices[0].message.content
    return json.loads(c[c.find("["):c.rfind("]")+1])

def screen_upload():
    ppt = st.session_state.ppt_data
    st.markdown("""
<div class="ac-overline" style="margin-bottom:8px">Volume I</div>
<div class="ac-h1">Upload Your<br><em>Lecture Deck</em></div>
<div class="ac-sub">Present a <code style="font-family:var(--fb);background:var(--muted);padding:1px 6px;border-radius:2px">.pptx</code> — the scholar's engine shall compose your examination.</div>
""", unsafe_allow_html=True)

    c1, c2 = st.columns([2, 1], gap="large")
    with c1:
        if ppt is None:
            st.markdown("""<div class="ac-drop">
  <div class="ac-arch" aria-hidden="true">📜</div>
  <div style="font-family:var(--fh);font-size:22px;font-weight:400;color:var(--fg);margin-bottom:6px">Choose a file or drag it here</div>
  <div style="font-family:var(--fb);font-size:14px;color:var(--mfg)">.pptx manuscripts only</div>
</div>""", unsafe_allow_html=True)
        up = st.file_uploader("Upload", type=["pptx"], label_visibility="collapsed")
        if up and ppt is None:
            with st.spinner(""):
                st.session_state.ppt_data = parse_pptx(up)
            st.rerun()
        if ppt:
            st.markdown(f"""<div class="ac-card ac-anim" style="margin-top:16px">
  <div class="ac-overline" style="margin-bottom:12px">Manuscript received</div>
  <div style="font-family:var(--fh);font-size:26px;font-weight:400;color:var(--fg);margin-bottom:18px">{ppt['filename']}</div>
  <div class="ac-div"></div>
  <div style="display:flex;gap:28px;margin-top:18px">
    <div>
      <div class="ac-overline" style="margin-bottom:4px">Folios</div>
      <div style="font-family:var(--fh);font-size:32px;color:var(--ac)">{ppt['slides']}</div>
    </div>
    <div>
      <div class="ac-overline" style="margin-bottom:4px">Words</div>
      <div style="font-family:var(--fh);font-size:32px;color:var(--ac)">{ppt['words']:,}</div>
    </div>
    <div>
      <div class="ac-overline" style="margin-bottom:4px">Status</div>
      <div style="font-family:var(--fh);font-size:32px;color:#7AAF95">&#x2713;</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Proceed to Configuration →", use_container_width=True):
                st.session_state.screen = "configure"; st.rerun()

    with c2:
        st.markdown("""<div class="ac-inset">
  <div class="ac-overline" style="margin-bottom:16px">The Method</div>
  <div style="display:flex;flex-direction:column;gap:18px">
    <div style="display:flex;gap:14px;align-items:flex-start">
      <div style="font-family:var(--fd);font-size:11px;color:var(--ac);padding-top:2px;flex-shrink:0">I</div>
      <div style="font-family:var(--fb);font-size:15px;color:var(--mfg);line-height:1.55">Upload your <em>.pptx</em> — every text shape is extracted and catalogued</div>
    </div>
    <div style="display:flex;gap:14px;align-items:flex-start">
      <div style="font-family:var(--fd);font-size:11px;color:var(--ac);padding-top:2px;flex-shrink:0">II</div>
      <div style="font-family:var(--fb);font-size:15px;color:var(--mfg);line-height:1.55">Set the difficulty, question count, and examination mode</div>
    </div>
    <div style="display:flex;gap:14px;align-items:flex-start">
      <div style="font-family:var(--fd);font-size:11px;color:var(--ac);padding-top:2px;flex-shrink:0">III</div>
      <div style="font-family:var(--fb);font-size:15px;color:var(--mfg);line-height:1.55">The AI composes MCQs with detailed per-option elucidations</div>
    </div>
  </div>
</div>""", unsafe_allow_html=True)

def screen_configure():
    ppt = st.session_state.ppt_data
    cfg = st.session_state.quiz_config
    diff = cfg.get("difficulty", "Medium")
    is_timed = cfg.get("timed", False)

    st.markdown(f'<div class="ac-chip">📎 &nbsp;{ppt["filename"]}</div>', unsafe_allow_html=True)
    st.markdown("""
<div class="ac-overline" style="margin-bottom:8px">Volume II</div>
<div class="ac-h1">Configure<br><em>Your Examination</em></div>
<div class="ac-sub">Set the breadth, rigour, and conditions of the test.</div>
""", unsafe_allow_html=True)

    c1, c2 = st.columns([2, 1], gap="large")
    with c1:
        st.markdown('<div class="ac-overline" style="margin-bottom:12px">Number of Questions</div>', unsafe_allow_html=True)
        num_q = st.slider("", 5, 30, cfg.get("num_questions", 10), key="cfg_n", label_visibility="collapsed")

        # ── Difficulty ──
        st.markdown('<div class="ac-overline" style="margin-top:24px;margin-bottom:12px">Difficulty</div>', unsafe_allow_html=True)
        dc = st.columns(3)
        for i, d in enumerate(["Simple", "Medium", "Complex"]):
            with dc[i]:
                with st.container(key=f"cfg_diff_{d}"):
                    if st.button(d, key=f"d_{d}", use_container_width=True):
                        st.session_state.quiz_config["difficulty"] = d; st.rerun()
                if d == diff:
                    st.markdown(f"""<style>
div[data-testid="cfg_diff_{d}"] button{{
  background:var(--bg-alt)!important;
  border:2px solid var(--ac)!important;
  color:var(--ac)!important;
  box-shadow:0 0 0 0!important;
  filter:none!important;
}}
</style>""", unsafe_allow_html=True)

        # ── Mode ──
        st.markdown('<div class="ac-overline" style="margin-top:24px;margin-bottom:12px">Examination Mode</div>', unsafe_allow_html=True)
        tc = st.columns(2)
        for mode_label, mode_val, key in [("⏱  Timed", True, "cfg_mode_timed"), ("∞  Untimed", False, "cfg_mode_untimed")]:
            with tc[0 if mode_val else 1]:
                with st.container(key=key):
                    if st.button(mode_label, key=f"m_{key}", use_container_width=True):
                        st.session_state.quiz_config["timed"] = mode_val; st.rerun()
                if is_timed == mode_val:
                    st.markdown(f"""<style>
div[data-testid="{key}"] button{{
  background:var(--bg-alt)!important;
  border:2px solid var(--ac)!important;
  color:var(--ac)!important;
  box-shadow:0 0 0 0!important;
  filter:none!important;
}}
</style>""", unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("Compose the Examination →", use_container_width=True):
            fd = cfg.get("difficulty", "Medium")
            ft = cfg.get("timed", False)
            st.session_state.quiz_config.update({
                "num_questions": num_q, "difficulty": fd, "timed": ft,
                "time_per_question": _TIME_PER_Q[fd] if ft else None,
            })
            st.session_state.current_question = 0
            st.session_state.user_answers = {}
            st.session_state.confirmed_answers = {}
            st.session_state.timer_start = time.time() if ft else None

            ph = st.empty()
            msgs = [
                f"Consulting {ppt['slides']} folios…",
                "Analysing scholarly content…",
                f"Composing question I of {num_q}…",
                f"Composing question {num_q//2} of {num_q}…",
                "Inscribing elucidations…",
            ]
            for m in msgs:
                ph.markdown(f"""<div style="padding:28px 0">
  <div style="font-family:var(--fb);font-size:15px;color:var(--mfg);margin-bottom:12px;font-style:italic">{m}</div>
  <div class="ac-lbar"><div class="ac-lbarf"></div></div>
</div>""", unsafe_allow_html=True)
                time.sleep(0.4)
            qs = generate_questions(ppt, num_q, fd)
            st.session_state.questions = qs
            ph.empty()
            st.session_state.screen = "quiz"; st.rerun()

    with c2:
        fd = cfg.get("difficulty", "Medium")
        ft = cfg.get("timed", False)
        tpq = _TIME_PER_Q[fd]
        total_s = num_q * tpq
        m_t, s_t = divmod(total_s, 60)
        st.markdown(f"""<div class="ac-inset">
  <div class="ac-overline" style="margin-bottom:14px">Duration Estimate</div>
  <div style="font-family:var(--fh);font-size:40px;font-weight:400;color:var(--ac);line-height:1">
    {m_t:02d}<span style="font-family:var(--fd);font-size:13px;color:var(--mfg)"> min </span>{s_t:02d}<span style="font-family:var(--fd);font-size:13px;color:var(--mfg)"> sec</span>
  </div>
  <div style="font-family:var(--fb);font-size:13px;color:var(--mfg);margin-top:6px">{'Timed' if ft else 'Untimed'} &middot; {tpq}s per question</div>
  <div class="ac-div" style="margin:18px 0"></div>
  <div style="font-family:var(--fb);font-size:14px;color:var(--mfg);line-height:2">
    Simple &nbsp;<span style="font-family:var(--fh);color:var(--fg)">~ 35 s/q</span><br>
    Medium &nbsp;<span style="font-family:var(--fh);color:var(--fg)">~ 67 s/q</span><br>
    Complex <span style="font-family:var(--fh);color:var(--fg)">~ 105 s/q</span>
  </div>
</div>""", unsafe_allow_html=True)

def screen_quiz():
    qs = st.session_state.questions
    qi = st.session_state.current_question
    q = qs[qi]
    cfg = st.session_state.quiz_config
    total = len(qs)
    sel = st.session_state.user_answers.get(q["id"])
    pct = (qi + 1) / total * 100

    # Overlay trick for invisible clickable buttons on styled option divs
    st.markdown("""<style>
div[data-testid^="opt_btn_"] .qz-opt { margin-bottom:0!important; }
</style>""", unsafe_allow_html=True)

    # Progress bar + counter
    st.markdown(f"""<div style="display:flex;align-items:center;gap:16px;margin-bottom:28px">
  <div class="ac-prog" style="flex:1"><div class="ac-progf" style="width:{pct}%"></div></div>
  <div style="font-family:var(--fd);font-size:10px;font-weight:500;color:var(--mfg);white-space:nowrap;letter-spacing:.15em">{qi+1:02d} / {total:02d}</div>
</div>""", unsafe_allow_html=True)

    c1, c2 = st.columns([2, 1], gap="large")
    with c1:
        # Question card
        st.markdown(f"""<div class="ac-card" style="margin-bottom:18px">
  <div class="ac-overline" style="margin-bottom:10px">Question {qi+1}</div>
  <div style="font-family:var(--fh);font-size:24px;font-weight:400;line-height:1.3;color:var(--fg)">{q['question_text']}</div>
</div>""", unsafe_allow_html=True)

        # Answer options
        for k, txt in q["options"].items():
            cls = "sel" if k == sel else ""
            with st.container(key=f"opt_btn_{qi}_{k}"):
                st.markdown(f'<div class="ac-opt {cls}"><span class="ac-key">{k}</span>{txt}</div>', unsafe_allow_html=True)
                if st.button(f"{k}. {txt}", key=f"o_{qi}_{k}", use_container_width=True):
                    st.session_state.user_answers[q["id"]] = k; st.rerun()

        st.markdown("<br>", unsafe_allow_html=True)
        nc1, nc2 = st.columns(2)
        with nc1:
            if qi > 0 and st.button("← Back", key=f"bk_{qi}", use_container_width=True):
                st.session_state.current_question -= 1; st.rerun()
        with nc2:
            if qi < total - 1:
                if st.button("Next →", key=f"nx_{qi}", use_container_width=True):
                    st.session_state.current_question += 1; st.rerun()
            else:
                if st.button("View Results →", key="res", use_container_width=True):
                    st.session_state.screen = "results"; st.rerun()

    with c2:
        # Timer
        if cfg.get("timed") and st.session_state.timer_start:
            elapsed = int(time.time() - st.session_state.timer_start)
            tpq = cfg["time_per_question"]
            rem = max(0, tpq - (elapsed % tpq))
            frac = rem / tpq
            urg = frac < 0.10
            R = 30; circ = 2 * math.pi * R; off = circ * (1 - frac)
            stk = "#B87070" if urg else "var(--ac)"
            st.markdown(f"""<div class="ac-inset" style="text-align:center;margin-bottom:14px">
  <div class="ac-overline" style="text-align:center;margin-bottom:12px">Time Remaining</div>
  <div class="ac-tw">
    <svg style="transform:rotate(-90deg)" width="72" height="72" viewBox="0 0 72 72">
      <circle fill="none" stroke="var(--muted)" stroke-width="3" cx="36" cy="36" r="{R}"/>
      <circle fill="none" stroke="{stk}" stroke-width="3" stroke-linecap="round"
        cx="36" cy="36" r="{R}" stroke-dasharray="{circ:.1f}" stroke-dashoffset="{off:.1f}"
        style="transition:stroke-dashoffset .9s linear,stroke .3s ease"/>
    </svg>
    <div class="ac-tn {'urg' if urg else ''}">{rem}</div>
  </div>
  <div style="font-family:var(--fd);font-size:9px;letter-spacing:.2em;color:var(--mfg)">SECONDS</div>
</div>""", unsafe_allow_html=True)
            st.rerun()

        # Question nav grid
        answered = st.session_state.user_answers
        st.markdown('<div class="ac-inset"><div class="ac-overline" style="margin-bottom:12px">Questions</div>', unsafe_allow_html=True)
        cols_per_row = 5
        for row_start in range(0, total, cols_per_row):
            row_qs = list(range(row_start, min(row_start + cols_per_row, total)))
            row_cols = st.columns(cols_per_row)
            for col_i in range(cols_per_row):
                with row_cols[col_i]:
                    if col_i < len(row_qs):
                        i = row_qs[col_i]
                        is_cur = i == qi
                        is_ans = qs[i]["id"] in answered
                        with st.container(key=f"qngrid_{i}"):
                            if st.button(str(i + 1), key=f"qn_{i}", use_container_width=True):
                                st.session_state.current_question = i; st.rerun()
                        if is_cur:
                            st.markdown(f'<style>div[data-testid="qngrid_{i}"] button{{background:var(--brass-grad)!important;color:var(--ac-fg)!important;border:none!important;}}</style>', unsafe_allow_html=True)
                        elif is_ans:
                            st.markdown(f'<style>div[data-testid="qngrid_{i}"] button{{background:rgba(201,169,98,.12)!important;color:var(--ac)!important;border-color:rgba(201,169,98,.4)!important;}}</style>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

        # Source slide
        src = q.get("source_slide")
        if src:
            st.markdown(f"""<div class="ac-inset" style="margin-top:12px">
  <div class="ac-overline" style="margin-bottom:8px">Source</div>
  <div style="font-family:var(--fh);font-size:28px;color:var(--fg)">Folio <span style="color:var(--ac)">{src:02d}</span></div>
</div>""", unsafe_allow_html=True)

def screen_results():
    qs = st.session_state.questions
    ans = st.session_state.user_answers
    right = sum(1 for q in qs if ans.get(q["id"]) == q["correct_answer"])
    total = len(qs)
    pct = int(right / total * 100) if total else 0
    cfg = st.session_state.quiz_config
    grade = "Distinction" if pct >= 80 else ("Credit" if pct >= 60 else "Refer for Re-examination")

    st.markdown("""
<div class="ac-overline" style="margin-bottom:8px">Volume IV</div>
<div class="ac-h1">Examination<br><em>Results</em></div>
""", unsafe_allow_html=True)

    c1, c2 = st.columns([2, 1], gap="large")
    with c1:
        # Score card with wax seal
        st.markdown(f"""<div class="ac-card ac-anim" style="margin-bottom:28px">
  <div style="display:flex;align-items:flex-start;justify-content:space-between">
    <div>
      <div class="ac-overline" style="margin-bottom:12px">Final Score</div>
      <div class="ac-score">{right}<span class="ac-score-denom"> / {total}</span></div>
      <div style="font-family:var(--fh);font-size:28px;color:var(--mfg);margin-top:8px">{pct}%</div>
    </div>
    <div style="text-align:right">
      <div class="ac-seal" title="{grade}" style="width:56px;height:56px;font-size:22px">
        {'★' if pct >= 80 else ('◆' if pct >= 60 else '◇')}
      </div>
      <div style="font-family:var(--fd);font-size:9px;color:var(--mfg);letter-spacing:.2em;margin-top:8px;text-transform:uppercase">{grade}</div>
    </div>
  </div>
  <div class="ac-div" style="margin:20px 0"></div>
  <div style="font-family:var(--fb);font-size:15px;color:var(--mfg);font-style:italic">
    {'A commendable performance — the scholar has demonstrated mastery of the subject.' if pct >= 80
     else ('A creditable result — further study of the weaker areas is advised.' if pct >= 60
           else 'The examination reveals gaps in understanding. A thorough review of the material is recommended.')}
  </div>
</div>""", unsafe_allow_html=True)

        # Question review
        st.markdown('<div class="ac-overline" style="margin-bottom:14px">Question Review</div>', unsafe_allow_html=True)
        for i, q in enumerate(qs):
            ua = ans.get(q["id"])
            ok = ua == q["correct_answer"]
            icon = "✓" if ok else "✗"
            with st.expander(f"{icon}  Q{i+1}: {q['question_text'][:72]}{'…' if len(q['question_text']) > 72 else ''}"):
                if ok:
                    st.markdown(f'<div class="ac-xrow xr-ok"><div class="ac-xlbl">Correct — {ua}: {q["options"].get(ua,"")}</div>{q["explanation_correct"]}</div>', unsafe_allow_html=True)
                else:
                    we = q["explanation_wrong"].get(ua, "") if ua else ""
                    st.markdown(f'<div class="ac-xrow xr-err" style="margin-bottom:6px"><div class="ac-xlbl">Your Answer — {ua}: {q["options"].get(ua, "—")}</div>{we}</div><div class="ac-xrow xr-ok"><div class="ac-xlbl">Correct Answer — {q["correct_answer"]}: {q["options"][q["correct_answer"]]}</div>{q["explanation_correct"]}</div>', unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)
        bc1, bc2 = st.columns(2)
        with bc1:
            if st.button("Retake Examination", use_container_width=True):
                st.session_state.current_question = 0
                st.session_state.user_answers = {}
                st.session_state.confirmed_answers = {}
                st.session_state.timer_start = time.time() if cfg.get("timed") else None
                st.session_state.screen = "quiz"; st.rerun()
        with bc2:
            if st.button("Upload New Manuscript →", use_container_width=True):
                for k, v in _DEFAULTS.items(): st.session_state[k] = v
                st.rerun()

    with c2:
        st.markdown(f"""<div class="ac-inset">
  <div class="ac-overline" style="margin-bottom:16px">Session Ledger</div>
  <div style="display:flex;flex-direction:column;gap:12px">
    <div style="display:flex;justify-content:space-between;align-items:baseline">
      <span style="font-family:var(--fb);font-size:14px;color:var(--mfg)">Difficulty</span>
      <span style="font-family:var(--fd);font-size:10px;color:var(--fg);letter-spacing:.1em">{cfg.get('difficulty','—')}</span>
    </div>
    <div style="display:flex;justify-content:space-between;align-items:baseline">
      <span style="font-family:var(--fb);font-size:14px;color:var(--mfg)">Mode</span>
      <span style="font-family:var(--fd);font-size:10px;color:var(--fg);letter-spacing:.1em">{'Timed' if cfg.get('timed') else 'Untimed'}</span>
    </div>
    <div style="display:flex;justify-content:space-between;align-items:baseline">
      <span style="font-family:var(--fb);font-size:14px;color:var(--mfg)">Questions</span>
      <span style="font-family:var(--fh);font-size:18px;color:var(--fg)">{total}</span>
    </div>
    <div style="display:flex;justify-content:space-between;align-items:baseline">
      <span style="font-family:var(--fb);font-size:14px;color:var(--mfg)">Correct</span>
      <span style="font-family:var(--fh);font-size:18px;color:var(--ac)">{right}</span>
    </div>
    <div class="ac-div"></div>
    <div style="display:flex;justify-content:space-between;align-items:baseline">
      <span style="font-family:var(--fb);font-size:14px;color:var(--mfg)">Verdict</span>
      <span style="font-family:var(--fd);font-size:9px;color:var(--ac);letter-spacing:.15em;text-transform:uppercase">{grade}</span>
    </div>
  </div>
</div>""", unsafe_allow_html=True)


def main():
    inject_css()
    scr = st.session_state.screen
    if scr == "upload":    screen_upload()
    elif scr == "configure": screen_configure()
    elif scr == "quiz":    screen_quiz()
    elif scr == "results": screen_results()

if __name__ == "__main__":
    main()
