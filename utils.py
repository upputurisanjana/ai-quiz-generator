import streamlit as st
from pptx import Presentation
import os, json, time
from dotenv import load_dotenv

load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), ".env"))

_TIME_PER_Q = {"Simple": 35, "Medium": 67, "Complex": 105}

_DEFAULTS = {
    "ppt_data": None,
    "quiz_config": {}, "questions": [], "current_question": 0,
    "user_answers": {}, "confirmed_answers": {}, "timer_start": None,
}

def init_state():
    for k, v in _DEFAULTS.items():
        if k not in st.session_state:
            st.session_state[k] = v

def inject_css():
    st.markdown("""
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,400;0,500;0,600;1,400;1,500&family=Crimson+Pro:ital,wght@0,400;0,600;1,400&family=Cinzel:wght@400;500;600&display=swap" rel="stylesheet">
<style>
:root {
  --bg:#F5F0E8;--bg-alt:#EDE7DB;--fg:#2C2218;--muted:#D8CFC3;--mfg:#7A6A58;
  --bd:#C8BDB0;--ac:#8B6914;--ac2:#8B2635;--ac-fg:#FFFFFF;
  --brass-grad:linear-gradient(180deg,#A07820 0%,#8B6914 50%,#7A5C10 100%);
  --fh:'Cormorant Garamond',serif;--fb:'Crimson Pro',serif;--fd:'Cinzel',serif;
}
html,body,.stApp{background:var(--bg)!important;color:var(--fg);font-family:var(--fb);}
.stApp>header,.stDeployButton,footer,.st-emotion-cache-18ni7ap,.st-emotion-cache-6qob1r{display:none!important;}
section[data-testid="stSidebar"]{display:none!important;}
.main .block-container{padding:52px 60px 100px!important;max-width:100%!important;min-height:100vh;}
.ac-texture{position:fixed;inset:0;pointer-events:none;z-index:0;background-image:url("data:image/svg+xml,%3Csvg xmlns=%22http://www.w3.org/2000/svg%22 width=%22300%22 height=%22300%22%3E%3Cfilter id=%22n%22%3E%3CfeTurbulence type=%22fractalNoise%22 baseFrequency=%220.75%22 numOctaves=%224%22 stitchTiles=%22stitch%22/%3E%3C/filter%3E%3Crect width=%22300%22 height=%22300%22 filter=%22url(%23n)%22 opacity=%220.03%22/%3E%3C/svg%3E");mix-blend-mode:overlay;}
.ac-vignette{position:fixed;inset:0;pointer-events:none;z-index:0;background:radial-gradient(ellipse at center,transparent 0%,transparent 50%,rgba(61,48,40,.5) 100%);}
.ac-overline{font-family:var(--fd);font-size:10px;font-weight:500;text-transform:uppercase;letter-spacing:.25em;color:var(--ac);}
.ac-h1{font-family:var(--fh);font-size:52px;font-weight:400;line-height:1.05;letter-spacing:-.01em;color:var(--fg);margin:8px 0 14px;}
.ac-sub{font-family:var(--fb);font-size:17px;color:var(--mfg);line-height:1.6;margin:0 0 36px;}
.ac-div{position:relative;height:1px;margin:24px 0;background:linear-gradient(90deg,transparent 0%,var(--bd) 20%,var(--ac) 50%,var(--bd) 80%,transparent 100%);}
.ac-div::before{content:"✶";position:absolute;left:50%;top:50%;transform:translate(-50%,-50%);color:var(--ac);font-size:11px;background:var(--bg-alt);padding:0 10px;}
.ac-card{background:var(--bg-alt);border:1px solid var(--bd);border-radius:4px;padding:28px 32px;position:relative;overflow:hidden;transition:border-color .3s,box-shadow .3s;}
.ac-card:hover{border-color:rgba(139,105,20,.45);box-shadow:0 8px 24px rgba(0,0,0,.1);}
.ac-card::before,.ac-card::after{content:"";position:absolute;width:24px;height:24px;border:1.5px solid var(--ac);opacity:.5;transition:opacity .3s;}
.ac-card::before{top:0;left:0;border-right:none;border-bottom:none;}
.ac-card::after{bottom:0;right:0;border-left:none;border-top:none;}
.ac-card:hover::before,.ac-card:hover::after{opacity:1;}
.ac-inset{background:var(--bg);border:1px solid var(--bd);border-radius:4px;padding:22px 24px;}
.ac-chip{display:inline-flex;align-items:center;gap:8px;background:rgba(139,105,20,.08);border:1px solid rgba(139,105,20,.25);border-radius:2px;padding:4px 14px;font-family:var(--fd);font-size:10px;font-weight:500;text-transform:uppercase;letter-spacing:.15em;color:var(--ac);margin-bottom:24px;}
.ac-drop{border:1px solid var(--bd);border-radius:4px;padding:56px 36px;text-align:center;background:var(--bg-alt);transition:border-color .3s;}
.ac-drop:hover{border-color:rgba(139,105,20,.45);}
.ac-arch{width:80px;height:80px;margin:0 auto 20px;border:1px solid var(--bd);border-radius:40% 40% 0 0/20% 20% 0 0;display:flex;align-items:center;justify-content:center;background:var(--bg);color:var(--ac);font-size:28px;transition:border-color .3s;}
.ac-drop:hover .ac-arch{border-color:rgba(139,105,20,.6);}
.ac-prog{height:2px;background:var(--muted);border-radius:99px;overflow:hidden;margin-bottom:28px;}
.ac-progf{height:100%;background:var(--brass-grad);border-radius:99px;transition:width .5s ease-out;}
.ac-lbar{width:100%;height:2px;background:var(--muted);border-radius:99px;overflow:hidden;}
.ac-lbarf{height:100%;background:var(--ac);border-radius:99px;animation:lsc 1.6s ease-in-out infinite;}
.ac-tw{position:relative;width:72px;height:72px;margin:0 auto 8px;}
.ac-tn{position:absolute;inset:0;display:flex;align-items:center;justify-content:center;font-family:var(--fh);font-size:20px;font-weight:500;color:var(--fg);transition:color .3s;}
.ac-tn.urg{color:#B87070;animation:pulse 1s ease infinite;}
.ac-score{font-family:var(--fh);font-size:72px;font-weight:400;letter-spacing:-.02em;line-height:1;color:var(--ac);}
.ac-score-denom{font-family:var(--fh);font-size:36px;color:var(--mfg);}
.ac-seal{display:inline-flex;align-items:center;justify-content:center;border-radius:50%;background:radial-gradient(circle at 40% 35%,#A83248 0%,#8B2635 55%,#6B1A28 100%);box-shadow:inset 0 2px 4px rgba(255,255,255,.2),inset 0 -2px 4px rgba(0,0,0,.3),0 3px 7px rgba(0,0,0,.4);color:#E8DFD4;flex-shrink:0;}
.ac-xrow{padding:14px 18px;font-family:var(--fb);font-size:15px;line-height:1.6;border-left:2px solid;margin-bottom:5px;border-radius:0 4px 4px 0;}
.xr-ok{border-color:#4A7C5F;background:rgba(74,124,95,.08);}
.xr-err{border-color:#8B3030;background:rgba(139,48,48,.08);}
.ac-xlbl{font-family:var(--fd);font-size:9px;font-weight:500;text-transform:uppercase;letter-spacing:.2em;margin-bottom:6px;}
.xr-ok .ac-xlbl{color:#7AAF95;}
.xr-err .ac-xlbl{color:#B87070;}
div.stButton>button[kind="secondary"]{background:var(--bg-alt)!important;color:var(--fg)!important;border:1px solid var(--bd)!important;border-radius:4px!important;font-family:var(--fd)!important;font-size:11px!important;font-weight:500!important;text-transform:uppercase!important;letter-spacing:.15em!important;padding:12px 24px!important;box-shadow:none!important;transition:border-color .15s,color .15s!important;}
div.stButton>button[kind="secondary"]:hover{border-color:var(--ac)!important;color:var(--ac)!important;}
div.stButton>button[kind="primary"]{background:rgba(139,105,20,.1)!important;color:var(--ac)!important;border:2px solid var(--ac)!important;border-radius:4px!important;font-family:var(--fd)!important;font-size:11px!important;font-weight:500!important;text-transform:uppercase!important;letter-spacing:.15em!important;padding:12px 24px!important;box-shadow:none!important;}
div.stButton>button[kind="primary"]:hover{background:rgba(139,105,20,.18)!important;}
div.stButton>button:focus-visible{outline:none!important;box-shadow:0 0 0 2px var(--ac)!important;}
div[data-testid="quiz_options"] button[kind="secondary"],
div[data-testid="quiz_options"] button[kind="primary"]{font-size:16px!important;text-transform:none!important;letter-spacing:0!important;text-align:left!important;justify-content:flex-start!important;min-height:52px!important;height:auto!important;padding:14px 18px!important;}
div[data-testid="stSlider"] label{font-family:var(--fd)!important;font-size:10px!important;text-transform:uppercase!important;letter-spacing:.2em!important;color:var(--mfg)!important;}
div[data-testid="stSlider"] [data-baseweb="slider"] [role="slider"]{background:var(--ac)!important;border-color:var(--ac)!important;}
div[data-testid="stSlider"] [data-baseweb="slider"] [data-testid="stThumbValue"]{font-family:var(--fd)!important;color:var(--ac)!important;font-size:11px!important;}
div[data-testid="stExpander"]{background:var(--bg-alt)!important;border:1px solid var(--bd)!important;border-radius:4px!important;margin-bottom:6px!important;}
div[data-testid="stExpander"] summary{font-family:var(--fh)!important;font-size:16px!important;color:var(--fg)!important;padding:14px 18px!important;}
@keyframes lsc{0%{width:0;margin-left:0;}50%{width:55%;margin-left:22%;}100%{width:0;margin-left:100%;}}
@keyframes pulse{0%,100%{opacity:1;}50%{opacity:.45;}}
@keyframes fup{from{opacity:0;transform:translateY(8px);}to{opacity:1;transform:translateY(0);}}
.ac-anim{animation:fup .3s ease-out;}
</style>
<div class="ac-texture" aria-hidden="true"></div>
<div class="ac-vignette" aria-hidden="true"></div>
""", unsafe_allow_html=True)

def parse_pptx(f):
    prs = Presentation(f)
    slides = [{"n": i+1, "t": " ".join(sh.text for sh in s.shapes if hasattr(sh, "text"))} for i, s in enumerate(prs.slides)]
    return {"filename": f.name, "slides": len(slides), "words": sum(len(s["t"].split()) for s in slides), "raw": slides}

def generate_questions(ppt, n, diff):
    from openai import OpenAI
    key = os.environ.get("OPENROUTER_API_KEY")
    if not key:
        st.error("OPENROUTER_API_KEY missing"); st.stop()
    client = OpenAI(api_key=key, base_url="https://openrouter.ai/api/v1")
    txt = "\n\n".join(f"Slide {s['n']}: {s['t']}" for s in ppt["raw"])
    prompt = f"""Generate {n} MCQs at {diff} difficulty.
Return ONLY a JSON array:
[{{"id":"q1","question_text":"...","options":{{"A":"...","B":"...","C":"...","D":"..."}},"correct_answer":"A","source_slide":1,"explanation_correct":"...","explanation_wrong":{{"B":"...","C":"...","D":"..."}}}}]
Content:\n{txt}"""
    r = client.chat.completions.create(model="deepseek/deepseek-chat", max_tokens=min(4096, 300 * n),
        messages=[{"role": "user", "content": prompt}])
    c = r.choices[0].message.content or ""
    start, end = c.find("["), c.rfind("]")
    if start == -1 or end == -1:
        st.error(f"Model returned invalid response:\n\n{c[:500]}"); st.stop()
    return json.loads(c[start:end+1])
